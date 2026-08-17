"""Build the LeakLens deck: the complete unlearning + quantization story.
Part I  the unlearning problem (thesis)  ->  Part II the quantization twist + the tool + its flow  ->
Part III results  ->  Part IV takeaways. Plain style: Times New Roman, grayscale, no em dashes."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

F = "figures"
FONT = "Times New Roman"
BLACK = RGBColor(0, 0, 0); INK = RGBColor(0x1A, 0x1A, 0x1A); GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); DARK = RGBColor(0x33, 0x33, 0x33)
ALT = RGBColor(0xF2, 0xF2, 0xF2); PALE = RGBColor(0xE7, 0xE7, 0xE7); HOT = RGBColor(0xC7, 0x3E, 0x1D)
TITLE, SUB, EYE, HEAD, BODY, SUBH, CAP = 34, 18, 13, 30, 16, 18, 12

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _nodash(t):
    t = str(t).replace(" — ", ", ").replace("—", ", ")
    return t.replace(" ,", ",").replace(",  ", ", ").replace(", .", ".")


def slide(): return prs.slides.add_slide(BLANK)


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def line(tf, text, size, color=INK, bold=False, italic=False, first=False, bullet=False,
         space=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = ("–  " + _nodash(text)) if bullet else _nodash(text)
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT
    return p


def header(s, eyebrow, title):
    tf = tbox(s, 0.7, 0.42, 12.0, 0.5); line(tf, eyebrow.upper(), EYE, GRAY, bold=True, first=True, space=2)
    L = len(_nodash(title)); hs = HEAD if L <= 52 else (25 if L <= 64 else 21)
    tf2 = tbox(s, 0.7, 0.82, 12.2, 0.72); line(tf2, title, hs, BLACK, bold=True, first=True)
    rule = s.shapes.add_shape(1, Inches(0.72), Inches(1.62), Inches(12.0), Pt(1.4))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLACK; rule.line.fill.background(); rule.shadow.inherit = False


def bullets(s, items, x=0.75, y=1.9, w=11.8, size=BODY, space=10):
    tf = tbox(s, x, y, w, SH.inches - y - 0.4)
    for i, it in enumerate(items):
        txt, kind = it if isinstance(it, tuple) else (it, "b")
        if kind == "h":
            line(tf, txt, SUBH, BLACK, bold=True, first=(i == 0), space=5)
        elif kind == "p":
            line(tf, txt, size, INK, first=(i == 0), space=space)
        else:
            line(tf, txt, size, INK, first=(i == 0), bullet=True, space=space)
    return tf


def caption(s, text, y=6.95, x=0.75, w=11.8):
    line(tbox(s, x, y, w, 0.45), text, CAP, GRAY, italic=True, first=True)


def picture(s, path, x=None, y=1.9, maxw=8.6, maxh=4.9):
    iw, ih = Image.open(f"{F}/{path}").size; ar = iw / ih; W = maxw; H = W / ar
    if H > maxh: H = maxh; W = H * ar
    left = int((SW - Inches(W)) / 2) if x is None else Inches(x)
    s.shapes.add_picture(f"{F}/{path}", left, Inches(y), width=Inches(W))


def fig_slide(eyebrow, title, path, cap=None, maxw=9.4, maxh=5.0, y=1.85):
    s = slide(); header(s, eyebrow, title); picture(s, path, y=y, maxw=maxw, maxh=maxh)
    if cap: caption(s, cap)
    return s


def divider(eyebrow, title, sub):
    s = slide()
    r = s.shapes.add_shape(1, 0, Inches(3.0), SW, Pt(2)); r.fill.solid(); r.fill.fore_color.rgb = BLACK
    r.line.fill.background(); r.shadow.inherit = False
    line(tbox(s, 1.0, 2.0, 11.3, 0.6), eyebrow.upper(), EYE, GRAY, bold=True, first=True)
    line(tbox(s, 1.0, 2.4, 11.3, 0.9), title, TITLE, BLACK, bold=True, first=True)
    line(tbox(s, 1.0, 3.2, 11.3, 0.8), sub, SUB, GRAY, italic=True, first=True)
    return s


def pipeline(s, steps, y=2.4, bh=1.35):
    n = len(steps); gap = 0.3; bw = (12.0 - gap * (n - 1)) / n; x = 0.7
    for i, (t, sub) in enumerate(steps):
        box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(bw), Inches(bh))
        box.fill.solid(); box.fill.fore_color.rgb = ALT; box.line.color.rgb = DARK; box.line.width = Pt(1.2)
        box.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        line(tf, t, 13, BLACK, bold=True, first=True, align=PP_ALIGN.CENTER, space=2)
        if sub: line(tf, sub, 9.5, GRAY, align=PP_ALIGN.CENTER, space=0)
        if i < n - 1:
            line(tbox(s, x + bw - 0.02, y + bh / 2 - 0.22, gap + 0.04, 0.44, MSO_ANCHOR.MIDDLE),
                 "→", 20, DARK, bold=True, first=True, align=PP_ALIGN.CENTER)
        x += bw + gap


# ============================ TITLE ============================
s = slide()
r = s.shapes.add_shape(1, 0, Inches(2.7), SW, Pt(2)); r.fill.solid(); r.fill.fore_color.rgb = BLACK
r.line.fill.background(); r.shadow.inherit = False
line(tbox(s, 1.0, 1.35, 11.3, 1.2), "LeakLens", TITLE + 8, BLACK, bold=True, first=True)
tf = tbox(s, 1.0, 2.95, 11.3, 2.0)
line(tf, "Does unlearning survive deployment compression?", SUB, INK, first=True, space=8)
line(tf, "Auditing whether quantization undoes machine unlearning, and localizing where it leaks, with the logit lens", SUB, INK, space=8)
line(tbox(s, 1.0, 6.5, 11.3, 0.6),
     "A portfolio project extending the TUM master's thesis on mechanistic interpretability of unlearning", CAP, GRAY, first=True)

# ============================ OUTLINE ============================
s = slide(); header(s, "Outline", "The complete story, in four parts")
bullets(s, [
    ("Part I. The unlearning problem (from the thesis)", "h"),
    "Output metrics cannot tell unlearning methods apart; the logit lens shows most unlearning suppresses rather than erases.",
    ("Part II. The quantization twist (LeakLens)", "h"),
    "Training and quantization are both lossy compression. Does deployment-time 4-bit quantization undo the deletion? The tool, its pipeline, and the recovery-depth metric.",
    ("Part III. Results", "h"),
    "The backend decides whether unlearning survives; the calibration-set attack (an honest null); the real-GPTQ cross-check.",
    ("Part IV. Takeaways", "h"),
    "What is firm, what is open, and what ships.",
], y=2.0, space=8)

# ==================================================================
# PART I  -  THE UNLEARNING PROBLEM
# ==================================================================
divider("Part I", "The unlearning problem", "what the thesis established: suppression is not erasure")

s = slide(); header(s, "1 · Machine unlearning and the verification gap", "did the model forget, or just hide it?")
bullets(s, [
    ("What unlearning promises", "h"),
    "Remove specific knowledge from a trained model, post hoc, without retraining: for privacy (right to be forgotten), copyright, and safety.",
    ("The hard question", "h"),
    "Unlearning is validated at the OUTPUT: ask the model and check the answer is gone. But a model can look forgotten on the outside while holding the knowledge inside.",
    ("Suppression vs erasure", "h"),
    "The thesis showed, on the Hubble benchmark, that most unlearning SUPPRESSES the fact at the output while its internal representation stays intact, and that a paraphrase, a few fine-tune steps, or a change of precision can bring it back.",
])

fig_slide("1 · The output is blind", "at the output, the four unlearning methods look identical (8B)",
          "infill_bars_8b.png", maxw=9.2, maxh=4.2,
          cap="Every unlearned method drops to the never-knew floor. One output verdict; no way to tell the methods, or true erasure, apart.")

fig_slide("1 · Reading inside: the logit lens", "the internal readout separates them: suppression is not erasure (8B)",
          "logitlens_rank_median_8b.png", maxw=10.0, maxh=5.0,
          cap="Rank-by-layer of the gold token. NPO / GradDiff / RMU bury it; IDK keeps it at the surface; none achieves internal erasure. Invisible at the output.")

fig_slide("1 · One fact, layer by layer", "the model that knows converges; the unlearned one never surfaces the answer",
          "lenstable_contrast_0.png", maxw=9.4, maxh=5.3,
          cap="Top-1 logit-lens prediction per layer. Suppression, not erasure, in a single table.")

fig_slide("1 · Confidence is not knowledge", "unlearned models are as confident in a wrong token as the knowing model is in the right one",
          "lenstable_confidence_contradiction.png", maxw=9.0, maxh=4.6,
          cap="A probability metric cannot separate confidently knows from confidently deflects. The output-metric failure, made visual.")

# ==================================================================
# PART II  -  THE QUANTIZATION TWIST
# ==================================================================
divider("Part II", "The quantization twist", "deployment compression can undo the deletion")

s = slide(); header(s, "2 · Two compressions", "training compresses, unlearning deletes, quantization compresses again")
bullets(s, [
    "Conklin et al. (2026) frame LLM training as LOSSY COMPRESSION: a model retains only what serves its objective, and different families compress differently.",
    ("The chain LeakLens studies", "h"),
    "1. Training compresses the data into weights.  2. Unlearning claims to surgically DELETE one retained piece.  3. Quantization applies a second, much cruder compression (4-bit) before deployment.",
    ("The question", "h"),
    "Unlearning is validated at FP16. Models ship at 4-bit. Does the second compression UNDO the first deletion, and if so, where in the network does the fact come back?",
])

s = slide(); header(s, "2 · Prior work and the gap", "the phenomenon is known; a representational, multi-backend audit is not")
rows = [("Zhang et al., ICLR 2025", "Quantization restores forgotten knowledge (21% at fp16 -> 83% at 4-bit). Behavioral + weight-space; benchmarks fail to detect it."),
        ("Two-failure-modes (2026)", "Logit-lens view of quantization failure generally; calls mechanistic analysis of quantized models preliminary and fragmented."),
        ("Mechanistic unlearning (2024)", "Interpretability used to DO unlearning more robustly; does not test quantization.")]
from pptx.util import Inches as _I
t = s.shapes.add_table(len(rows) + 1, 2, _I(0.75), _I(1.95), _I(11.8), _I(0.4 * (len(rows) + 1))).table
t.columns[0].width = _I(3.6); t.columns[1].width = _I(8.2)
hdr = [("prior work", "what it left open")] + rows
for i, (a, b) in enumerate(hdr):
    for j, val in enumerate((a, b)):
        c = t.cell(i, j); c.text = _nodash(val); c.vertical_anchor = MSO_ANCHOR.MIDDLE
        rr = c.text_frame.paragraphs[0].runs[0]; rr.font.name = FONT
        if i == 0:
            c.fill.solid(); c.fill.fore_color.rgb = DARK; rr.font.bold = True; rr.font.color.rgb = WHITE; rr.font.size = Pt(13)
        else:
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else ALT; rr.font.size = Pt(12.5); rr.font.color.rgb = INK; rr.font.bold = (j == 0)
line(tbox(s, 0.75, 5.55, 11.8, 1.2),
     "The gap: no reproducible, multi-backend tool that reports REPRESENTATIONAL recovery and localizes the layer at which forgotten knowledge re-opens. LeakLens builds one with the logit lens.",
     BODY, BLACK, bold=True, first=True)

s = slide(); header(s, "2 · LeakLens: what it does", "a config-driven audit of unlearning under quantization")
bullets(s, [
    ("Inputs (one YAML)", "h"),
    "A base model that KNOWS the facts, the UNLEARNED models (one per method), a forget set, and a list of deployment quantization configs.",
    ("Five quantization backends", "h"),
    "none (fp16 reference), bitsandbytes (INT8 / NF4 / FP4), round-to-nearest (INT8 / INT4), calibration-aware AWQ, and real GPTQ.",
    ("Outputs", "h"),
    "Behavioral recovery, the representational RECOVERY-DEPTH metric, and a single heatmap (method x quant, colour = recovery, label = depth) plus a report card.",
])

s = slide(); header(s, "2 · How it works: the pipeline", "how execution travels through the tool (flow.md)")
pipeline(s, [("leaklens audit\n--config", "one YAML"),
             ("load base\nonce", "knows reference"),
             ("for method x\nquant", "load + quantize"),
             ("logit-lens\nprobe", "rank per layer"),
             ("recovery\nmetrics", "fraction + depth"),
             ("report", "heatmap")], y=2.5, bh=1.5)
bullets(s, [
    "The backend is a pluggable dimension: the same audit loop drives fp16, bitsandbytes, RTN, AWQ, and GPTQ. Calibration-based backends take a calibration corpus, and sweeping that corpus IS the calibration-set attack.",
    "Two compute-heavy leaves (the lens and the metrics) have no internal dependencies, so they are unit-tested in isolation; every number traces to one config file.",
], y=4.5, space=8)

s = slide(); header(s, "2 · The recovery-depth metric", "not just whether forgotten knowledge returns, but where")
bullets(s, [
    "The logit lens gives, at every layer, the rank of the gold answer token. Run it on the base, the unlearned, and the unlearned-then-quantized model.",
    ("Recovery depth (the novel scalar)", "h"),
    "The shallowest mid-stack layer at which the QUANTIZED model's gold-token rank returns toward the base while the fp16-unlearned model's had not, i.e. the layer at which quantization re-opens the fact.",
    ("Recovery fraction (the behavioral screen)", "h"),
    "The share of the deleted knowledge that returns at a given precision, normalized against the fp16-unlearned baseline. Isolates the quantization effect; drives the heatmap colour.",
], space=9)

# ==================================================================
# PART III  -  RESULTS
# ==================================================================
divider("Part III", "Results", "on the Hubble-8B suite (base recalls the facts at gold-probability 0.909)")

fig_slide("3 · Result 1: the backend decides survival", "one hot cell: NPO's unlearning is undone by 4-bit round-to-nearest",
          "heatmap_tier1_8b.png", maxw=8.8, maxh=4.6,
          cap="NPO's forgotten fact goes from rank 4,263 (fp16) to rank 3 under rtn-int4 (recovery 0.16), while calibration-aware bitsandbytes NF4 / INT8 keep it buried. The backend, not just the bit-width, decides.")

fig_slide("3 · Result 1: the recovery trajectory", "where the fact re-opens, layer by layer",
          "leaklens_trajectory_8b.png", maxw=8.6, maxh=4.8,
          cap="Base (knows) vs NPO-unlearned (fp16, buried) vs NPO + rtn-int4 (recovered). Quantization pulls the trajectory back toward the knowing model.")

s = slide(); header(s, "3 · Result 2: four mechanisms, four responses", "the internal readout explains who survives quantization")
bullets(s, [
    ("IDK", "h"),
    "Redirects the OUTPUT (\"I don't know\") without disturbing the representation, so the token stays internally accessible throughout; quantization neither helps nor hurts.",
    ("NPO", "h"),
    "Buries the fact deeply at fp16 but the representation is fragile: crude 4-bit RTN re-opens it (Result 1).",
    ("RMU", "h"),
    "Scrambles the representation early and deep (the strongest suppression of the three); only mildly moved by quantization.",
    "The deployment-compression view reproduces the thesis mechanistic taxonomy (IDK surface, NPO late suppression, RMU deep scramble).",
], space=6)

fig_slide("3 · Result 3: the calibration-set attack", "an honest null: calibration proximity barely moves recovery",
          "calibration_attack_8b.png", maxw=8.2, maxh=4.5,
          cap="Sweeping the AWQ calibration corpus generic -> adjacent -> forget moves recovery only 0.006 to 0.009, within noise at n=25. Not shown to be an attack surface. Reported as a null, not spun.")

s = slide(); header(s, "3 · Result 4: the real-GPTQ cross-check", "the decisive test could not be run in this environment")
bullets(s, [
    "We wired a real GPTQ backend (gptqmodel, desc_act, calibration = the swept corpus) as the decisive, more calibration-sensitive test.",
    ("What happened", "h"),
    "GPTQ QUANTIZES successfully, but the packed model's forward pass crashes in gptqmodel's fallback kernel, an incompatibility with this environment's bleeding-edge Transformers 5.14 / Torch 2.11 (no optimized CUDA kernels available). Three successive fixes each advanced further; the last failure is internal to the library.",
    ("Consequence, stated plainly", "h"),
    "The GPTQ cross-check is unresolved here, so the calibration-attack question stays open. The firm result LeakLens ships is the backend-dependent recovery of Result 1.",
], space=8)

# ==================================================================
# PART IV  -  TAKEAWAYS
# ==================================================================
divider("Part IV", "Takeaways", "the complete story, and what is firm vs open")

s = slide(); header(s, "4 · Limitations", "stated honestly")
bullets(s, [
    "Single model family (Hubble), single attribute (birthdate), one unlearned checkpoint per method; the recovery findings are on 25 to 44 facts.",
    "The AWQ backend is a transparent activation-scaling re-implementation isolating the calibration mechanism, not the reference kernel.",
    "The real-GPTQ cross-check was blocked by an environment/kernel incompatibility, so the calibration-set attack remains an open question rather than a settled result.",
    "Recovery depth and recovery fraction are complementary; the behavioral fraction is the primary, quantization-isolating signal.",
])

s = slide(); header(s, "Conclusion", "one continuous story, from unlearning to deployment")
tf = tbox(s, 1.0, 2.0, 11.3, 3.8, anchor=MSO_ANCHOR.MIDDLE)
line(tf, "The thesis showed that most unlearning SUPPRESSES a fact at the output while its", 20, BLACK, first=True, space=8, align=PP_ALIGN.CENTER)
line(tf, "representation stays intact. LeakLens shows that deployment-time quantization can", 20, BLACK, space=8, align=PP_ALIGN.CENTER)
line(tf, "RE-OPEN that suppressed representation, and that the quantization backend decides whether it does.", 20, BLACK, bold=True, space=8, align=PP_ALIGN.CENTER)
line(tf, "If your compliance story depends on unlearning, quantization can silently void it, and now there is a tool that measures it.", BODY, GRAY, italic=True, space=10, align=PP_ALIGN.CENTER)
line(tf, "Future work: a pinned GPTQ / AWQ environment to settle the calibration attack, more model families and attributes, and the PII per-attribute breakdown.", BODY, GRAY, italic=True, space=0, align=PP_ALIGN.CENTER)

out = "LeakLens.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
